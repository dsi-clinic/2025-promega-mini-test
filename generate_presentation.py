#!/usr/bin/env python3
"""Generate the Image Classifier presentation as a PPTX file.
Charts read from the newest overlay threshold study results CSV.
"""

import os
import glob
import tempfile
from io import BytesIO

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import numpy as np
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Colour palette (clean light theme)
# ---------------------------------------------------------------------------
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF7, 0xF7, 0xF7)
DARK_TEXT   = RGBColor(0x2D, 0x2D, 0x2D)
MID_GRAY    = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT_BLUE = RGBColor(0x1F, 0x77, 0xB4)
ACCENT_ORANGE = RGBColor(0xFF, 0x7F, 0x0E)
ACCENT_RED  = RGBColor(0xD6, 0x27, 0x28)
ACCENT_GREEN = RGBColor(0x2C, 0xA0, 0x2C)

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def set_slide_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text, left=Inches(0.6), top=Inches(0.3),
              width=Inches(12), height=Inches(0.8),
              font_size=Pt(32), bold=True, color=DARK_TEXT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    return txBox


def add_subtitle(slide, text, left=Inches(0.6), top=Inches(1.1),
                 width=Inches(12), height=Inches(0.5),
                 font_size=Pt(18), color=MID_GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.name = "Calibri"
    return txBox


def add_body_text(slide, bullets, left=Inches(0.6), top=Inches(1.7),
                  width=Inches(5.8), height=Inches(5.2),
                  font_size=Pt(16), color=DARK_TEXT, line_spacing=Pt(24)):
    """Add a list of bullet strings. Supports (text, indent_level) tuples."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        if line_spacing:
            p.line_spacing = line_spacing
    return txBox


def add_speaker_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def add_section_divider(slide, number, title):
    """Full-slide section divider."""
    set_slide_bg(slide, OFF_WHITE)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Section {number}"
    p.font.size = Pt(20)
    p.font.color.rgb = MID_GRAY
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = DARK_TEXT
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER


def add_accent_line(slide, left=Inches(0.6), top=Inches(1.05),
                    width=Inches(2), color=ACCENT_BLUE):
    """Thin accent line under the title."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def fig_to_image_stream(fig, dpi=180):
    """Convert a matplotlib figure to a BytesIO PNG stream."""
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight",
                facecolor="white", edgecolor="none")
    plt.close(fig)
    buf.seek(0)
    return buf


# ---------------------------------------------------------------------------
# Data loading (newest run)
# ---------------------------------------------------------------------------

def _find_newest_overlay_csv():
    """Path to the most recently modified overlay_threshold_study_results.csv."""
    root = os.path.dirname(os.path.abspath(__file__))
    candidates = [
        os.path.join(root, "2025-promega-mini-test", "comparison_runs", "overlay_threshold_study_results.csv"),
        os.path.join(root, "comparison_runs", "overlay_threshold_study_results.csv"),
    ]
    candidates += glob.glob(os.path.join(root, "regeneration", "run_*", "overlay_threshold_study_results.csv"))
    candidates = [p for p in candidates if os.path.isfile(p)]
    if not candidates:
        return None
    return max(candidates, key=os.path.getmtime)


def load_overlay_results(csv_path=None):
    """Load per_day and effnet_ts balanced_acc and optimal_threshold at optimal threshold per day.
    Returns (days, per_day_acc, effnet_acc, per_day_thresh, effnet_thresh).
    """
    import pandas as pd
    path = csv_path or _find_newest_overlay_csv()
    if path is None:
        # Fallback to original hardcoded values
        days = [6, 8, 10, 13, 15, 17, 20.5, 24, 26, 28, 30]
        return days, [0.663, 0.502, 0.616, 0.676, 0.619, 0.567, 0.610, 0.617, 0.617, 0.602, 0.629], [0.651, 0.500, 0.583, 0.449, 0.465, 0.538, 0.492, 0.594, 0.579, 0.816, 0.763], [0.51, 0.53, 0.55, 0.54, 0.56, 0.53, 0.52, 0.58, 0.58, 0.56, 0.58], [0.32, 0.50, 0.50, 0.22, 0.30, 0.89, 0.86, 0.16, 0.10, 0.90, 0.10]
    df = pd.read_csv(path)
    df["day"] = df["day"].astype(float)
    days_order = [6, 8, 10, 13, 15, 17, 20.5, 24, 26, 28, 30]
    df["th_round"] = df["threshold"].round(2)
    df["opt_round"] = df["optimal_threshold"].round(2)
    opt = df[df["th_round"] == df["opt_round"]].drop_duplicates(subset=["model", "day"], keep="first")
    per_day_df = opt[opt.model == "per_day"].set_index("day").reindex(days_order)
    effnet_df = opt[opt.model == "effnet_ts"].set_index("day").reindex(days_order)
    per_day_acc = [round(x, 3) for x in per_day_df["balanced_acc"].tolist()]
    effnet_acc = [round(x, 3) for x in effnet_df["balanced_acc"].tolist()]
    per_day_t = [round(x, 2) for x in per_day_df["optimal_threshold"].tolist()]
    effnet_t = [round(x, 2) for x in effnet_df["optimal_threshold"].tolist()]
    return days_order, per_day_acc, effnet_acc, per_day_t, effnet_t


# ---------------------------------------------------------------------------
# Chart generators
# ---------------------------------------------------------------------------

def make_performance_chart():
    """Line chart: per-day vs effnet_ts balanced accuracy across days (from newest run CSV)."""
    days, per_day, effnet = load_overlay_results()[:3]

    fig, ax = plt.subplots(figsize=(9, 4.5))
    x = np.arange(len(days))

    ax.plot(x, per_day, "o-", color="#1f77b4", linewidth=2.5, markersize=8,
            label="Per-Day EfficientNet", zorder=3)
    ax.plot(x, effnet, "s--", color="#ff7f0e", linewidth=2.5, markersize=8,
            label="EfficientNet Time-Series", zorder=3)

    # Highlight the crossover region
    ax.axvspan(8.5, 10.5, alpha=0.08, color="green", label="_nolegend_")

    # Mark day 13/15 collapse (use effnet value at index 3 = day 13)
    collapse_y = effnet[3] if len(effnet) > 3 else 0.45
    ax.annotate("Collapse", xy=(3, collapse_y), xytext=(3.8, min(0.38, collapse_y - 0.07)),
                fontsize=10, color="#d62728", fontweight="bold",
                arrowprops=dict(arrowstyle="->", color="#d62728", lw=1.5))

    # Mark crossover (effnet wins late days)
    late_idx = 9
    if len(effnet) > late_idx and effnet[late_idx] > per_day[late_idx]:
        ax.annotate("Crossover\n(Day 28+)", xy=(late_idx, effnet[late_idx]), xytext=(7.5, 0.87),
                    fontsize=10, color="#2ca02c", fontweight="bold",
                    arrowprops=dict(arrowstyle="->", color="#2ca02c", lw=1.5))

    # Shading: per-day wins vs effnet wins
    ax.fill_between(x, per_day, effnet, where=[p >= e for p, e in zip(per_day, effnet)],
                    alpha=0.07, color="#1f77b4")
    ax.fill_between(x, per_day, effnet, where=[e > p for p, e in zip(per_day, effnet)],
                    alpha=0.07, color="#ff7f0e")

    ax.set_xticks(x)
    ax.set_xticklabels([str(d) for d in days], fontsize=11)
    ax.set_xlabel("Day", fontsize=13, fontweight="bold")
    ax.set_ylabel("Balanced Accuracy", fontsize=13, fontweight="bold")
    ax.set_ylim(0.35, 0.95)
    ax.yaxis.set_major_formatter(mticker.FormatStrFormatter("%.2f"))
    ax.legend(fontsize=11, loc="upper left", framealpha=0.9)
    ax.grid(True, alpha=0.3)
    ax.set_title("Per-Day vs Time-Series: Balanced Accuracy by Day",
                 fontsize=14, fontweight="bold", pad=12)

    n_per_day_wins = sum(1 for p, e in zip(per_day, effnet) if p > e)
    ax.text(4, 0.91, f"Per-day wins {n_per_day_wins} of {len(days)} days",
            fontsize=11, color="#1f77b4", fontweight="bold",
            bbox=dict(boxstyle="round,pad=0.3", facecolor="white", edgecolor="#1f77b4", alpha=0.8))

    fig.tight_layout()
    return fig_to_image_stream(fig)


def make_threshold_chart():
    """Bar chart: optimal thresholds per day (from newest run CSV)."""
    days, _, _, per_day_t, effnet_t = load_overlay_results()

    fig, ax = plt.subplots(figsize=(9, 4.5))
    x = np.arange(len(days))
    w = 0.35

    bars1 = ax.bar(x - w/2, per_day_t, w, label="Per-Day (stable)", color="#1f77b4", alpha=0.85)
    bars2 = ax.bar(x + w/2, effnet_t, w, label="effnet_ts (unstable)", color="#ff7f0e", alpha=0.85)

    # Reference line at 0.5
    ax.axhline(y=0.5, color="#d62728", linestyle="--", linewidth=1.5, alpha=0.7, label="Default 0.5")

    ax.set_xticks(x)
    ax.set_xticklabels([str(d) for d in days], fontsize=11)
    ax.set_xlabel("Day", fontsize=13, fontweight="bold")
    ax.set_ylabel("Optimal Threshold", fontsize=13, fontweight="bold")
    ax.set_ylim(0, 1.05)
    ax.legend(fontsize=11, loc="upper right", framealpha=0.9)
    ax.grid(True, axis="y", alpha=0.3)
    ax.set_title("Optimal Threshold by Day: Per-Day (stable) vs Time-Series (wild swings)",
                 fontsize=13, fontweight="bold", pad=12)

    # Annotate extreme values
    min_t = min(effnet_t)
    max_t = max(effnet_t)
    if min_t <= 0.25:
        idx_min = effnet_t.index(min_t)
        ax.annotate(f"{min_t:.2f}", xy=(idx_min, min_t), xytext=(idx_min, 0.02),
                    fontsize=9, ha="center", color="#d62728", fontweight="bold")
    if max_t >= 0.85:
        idx_max = effnet_t.index(max_t)
        ax.annotate(f"{max_t:.2f}", xy=(idx_max, max_t), xytext=(idx_max, 0.97),
                    fontsize=9, ha="center", color="#d62728", fontweight="bold")

    fig.tight_layout()
    return fig_to_image_stream(fig)


def make_overlay_grid():
    """Create a 2x3 grid of overlay images: rows = Acceptable/Not Acceptable,
    cols = Day 3, Day 28, Day 30."""
    base = "/home/tonyluo/amanda_temporal/_tmp_imgs"
    paths = [
        # Row 1: Acceptable
        (f"{base}/day03_acc.png", "Day 3\nAcceptable"),
        (f"{base}/day28_acc.png", "Day 28\nAcceptable"),
        (f"{base}/day30_acc.png", "Day 30\nAcceptable"),
        # Row 2: Not Acceptable
        (f"{base}/day03_notacc.png", "Day 3\nNot Acceptable"),
        (f"{base}/day28_notacc.png", "Day 28\nNot Acceptable"),
        (f"{base}/day30_notacc.png", "Day 30\nNot Acceptable"),
    ]

    fig, axes = plt.subplots(2, 3, figsize=(10, 6))
    fig.suptitle("Overlay Images: Organoid Outline on Grayscale Background",
                 fontsize=14, fontweight="bold", y=1.02)

    for idx, (path, label) in enumerate(paths):
        row, col = divmod(idx, 3)
        ax = axes[row][col]
        try:
            img = Image.open(path)
            ax.imshow(np.array(img))
        except Exception:
            ax.text(0.5, 0.5, "Image\nnot found", ha="center", va="center",
                    transform=ax.transAxes, fontsize=12, color="gray")
        ax.set_title(label, fontsize=11, fontweight="bold", pad=4)
        ax.axis("off")

    # Row labels
    fig.text(0.02, 0.72, "Acceptable", fontsize=12, fontweight="bold",
             color="#2ca02c", rotation=90, va="center")
    fig.text(0.02, 0.30, "Not Acceptable", fontsize=12, fontweight="bold",
             color="#d62728", rotation=90, va="center")

    fig.tight_layout(rect=[0.04, 0, 1, 1])
    return fig_to_image_stream(fig, dpi=200)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]  # blank

    # =====================================================================
    # TITLE SLIDE (blank for user to fill in)
    # =====================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title(slide, "(Your Title Here)", top=Inches(2.5),
              left=Inches(2), width=Inches(9),
              font_size=Pt(40), color=MID_GRAY)
    add_subtitle(slide, "Add your name, date, and affiliation",
                 top=Inches(3.5), left=Inches(2), width=Inches(9),
                 font_size=Pt(20))
    add_speaker_notes(slide, "Replace this with your presentation title, name, and date.")

    # =====================================================================
    # SLIDE 1-1: Data Funnel and Dataset at a Glance
    # =====================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title(slide, "Data Funnel and Dataset at a Glance")
    add_accent_line(slide)

    bullets = [
        "The Data Funnel:",
        ("475 unique organoid wells imaged across 11 days", 1),
        ("\u2192 ~260 have Day 30 expert survey labels (majority vote, min 4 votes)", 1),
        ("\u2192 220 remain after requiring complete metabolite data + valid image + mask", 1),
        "",
        "Train / Val / Test Split (organoid-level):",
        ("158 train  /  18 validation  /  44 test organoids", 1),
        ("Test labels: 35 Acceptable, 9 Not Acceptable (per day)", 1),
        "",
        "11 Time Points: Day 3, 6, 8, 10, 13, 15, 17, 20.5, 24, 26, 28, 30",
        "",
        "Key takeaway: only 9 negative test samples per day",
    ]
    add_body_text(slide, bullets, width=Inches(11))

    add_speaker_notes(slide,
        "We start with 475 organoids, but after requiring survey labels and complete "
        "metabolite data, only 220 remain. The test set has just 9 negatives per day -- "
        "keep this number in mind; it comes back in Observation 3.")

    # =====================================================================
    # SLIDE 1-2: Labels, Class Imbalance, and Example Overlays
    # =====================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title(slide, "Labels, Class Imbalance, and What the Model Sees")
    add_accent_line(slide)

    bullets = [
        "Class split: 72.5% Acceptable vs 27.5% Not Acceptable",
        "Imbalance handled with balanced class weights:",
        ("\u2022 Minority class (~Not Acceptable) gets ~2.5\u00d7 weight in loss", 1),
        ("\u2022 sklearn compute_class_weight(\"balanced\")", 1),
        "",
        "Labels from expert surveys at Day 28/30,",
        "propagated backward to earlier days",
        "",
        "Image representation: overlay (green outline on grayscale) \u2192",
    ]
    add_body_text(slide, bullets, width=Inches(5.5))

    # Embed the overlay grid image on the right
    overlay_stream = make_overlay_grid()
    slide.shapes.add_picture(overlay_stream, Inches(6.4), Inches(1.5),
                             width=Inches(6.5))

    add_speaker_notes(slide,
        "The class imbalance is structural -- most organoids are acceptable. We up-weight "
        "the minority class ~2.5x during training, but 9 negatives in the test set still "
        "makes evaluation fragile. The overlay images show the organoid's boundary on a "
        "neutral background -- this is what the model actually sees. Notice how early-day "
        "organoids look similar regardless of label, while late-day ones show clear differences.")

    # =====================================================================
    # SLIDE 2-1: Pipeline and Image Representation (combined, brief)
    # =====================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title(slide, "Data Engineering Pipeline (Brief)")
    add_accent_line(slide)

    bullets = [
        "Pipeline in one line:",
        ("Raw TIFF \u2192 best-focus selection \u2192 resize 512\u00d7384 \u2192 SegFormer segmentation \u2192 overlay \u2192 mean-fill background \u2192 split", 1),
        "",
        "Key decisions:",
        ("\u2022 Two segmentation models: Early (Day 3-10) and Late (Day 13-30) due to morphology differences", 1),
        ("\u2022 Overlay representation = organoid shape on neutral background. Shape is the primary signal.", 1),
        ("\u2022 Other representations tested (RGB mask, filled mask) -- overlay performed best", 1),
        "",
        "This is a 2-stage pipeline: segmentation quality directly affects classification",
    ]
    add_body_text(slide, bullets, width=Inches(11.5), font_size=Pt(15))

    add_speaker_notes(slide,
        "15-step pipeline from raw TIFF to training-ready images. The key design decision "
        "is the overlay representation -- showing the organoid's boundary on a neutral "
        "background. This works well, suggesting morphological shape is the main signal "
        "for quality classification. Keep this brief -- ~90 seconds.")

    # =====================================================================
    # SLIDE 3-1: Two Models, One Question
    # =====================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title(slide, "Two Models, One Question")
    add_accent_line(slide)

    bullets = [
        "Per-Day EfficientNet:",
        ("\u2022 One model per day. Sees only the image at day D.", 1),
        ("\u2022 No history, no temporal context.", 1),
        "",
        "EfficientNet Time-Series (effnet_ts):",
        ("\u2022 Sees the full image sequence up to day D.", 1),
        ("\u2022 Temporal context from all earlier days.", 1),
        "",
        "Same backbone (EfficientNet-B0), same split, same class weighting.",
        "",
        "Primary metric: balanced accuracy at validation-optimal threshold.",
        ("Note: thresholds are chosen by the validation set. At a naive 0.5 cutoff, results are worse.", 1),
        "",
        "The question: does seeing history help?",
    ]
    add_body_text(slide, bullets, width=Inches(11.5))

    add_speaker_notes(slide,
        "The core comparison: a simple per-day model vs a temporal model that sees the full "
        "history. Same backbone, same data, same split. The only difference is temporal context. "
        "Thresholds are optimized on validation -- at 0.5, performance would be worse.")

    # =====================================================================
    # SLIDE 3-2: Performance Chart (KEY SLIDE)
    # =====================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title(slide, "Per-Day vs Time-Series: Performance Across Days")
    add_accent_line(slide)

    perf_stream = make_performance_chart()
    slide.shapes.add_picture(perf_stream, Inches(1.5), Inches(1.4),
                             width=Inches(10))

    add_speaker_notes(slide,
        "Per-day wins 9 out of 11 days. Time-series only overtakes at days 28-30. "
        "The dips at day 13 and 15 are where the time-series model completely collapses. "
        "This is a striking result -- the simpler model wins almost everywhere. "
        "This leads us to our three observations.")

    # =====================================================================
    # SLIDE 4-1: Observation 1 -- Threshold Sensitivity
    # =====================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title(slide, "Obs 1: Model Calibration is Unstable")
    add_accent_line(slide, color=ACCENT_RED)

    thresh_stream = make_threshold_chart()
    slide.shapes.add_picture(thresh_stream, Inches(0.5), Inches(1.4),
                             width=Inches(7.5))

    side_bullets = [
        "effnet_ts thresholds:",
        ("\u2022 Range: 0.10 \u2013 0.90 (9\u00d7 spread)", 1),
        ("\u2022 Cannot deploy with a single cutoff", 1),
        "",
        "Per-day thresholds:",
        ("\u2022 Range: 0.51 \u2013 0.58 (stable)", 1),
        "",
        "Why? The temporal model",
        "absorbs noise from earlier",
        "days, distorting its",
        "probability outputs.",
        "",
        "\u2192 Connects to Obs 2",
    ]
    add_body_text(slide, side_bullets, left=Inches(8.3), top=Inches(1.4),
                  width=Inches(4.5), font_size=Pt(14))

    add_speaker_notes(slide,
        "The time-series model's optimal threshold swings from 0.10 to 0.90 across days. "
        "You cannot deploy this with a single cutoff. The per-day model stays in a tight "
        "0.51-0.58 range. This instability is because the temporal model is absorbing noise "
        "from earlier days -- which connects directly to Observation 2.")

    # =====================================================================
    # SLIDE 4-2: Observation 1 -- Day 13/15 Collapse
    # =====================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title(slide, "Obs 1 (cont.): Day 13/15 \u2014 The Extreme Collapse")
    add_accent_line(slide, color=ACCENT_RED)

    bullets = [
        "effnet_ts on Day 13:",
        ("\u2022 Predicted probabilities collapse to ~0.21 (range 0.207 \u2013 0.222)", 1),
        ("\u2022 At threshold 0.5: predicts ALL as \"Not Acceptable\"", 1),
        ("\u2022 Result: 0 true positives out of 31 \u2192 balanced accuracy = 0.50", 1),
        "",
        "Remediation attempts (briefly):",
        ("\u2022 Grayscale normalization \u2192 still collapsed", 1),
        ("\u2022 Filled mask input \u2192 still collapsed", 1),
        ("\u2022 Transform auditing (Day 13 vs Day 28) \u2192 no obvious input bug", 1),
        "",
        "Per-day model on same days: balanced_acc = 0.676 (Day 13), 0.619 (Day 15)",
        ("\u2192 Per-day handles these days fine. The issue is in temporal modeling.", 1),
    ]
    add_body_text(slide, bullets, width=Inches(11.5))

    add_speaker_notes(slide,
        "Days 13 and 15 are the extreme case -- probabilities collapse to 0.21. We tried "
        "three fixes, none worked. The per-day model handles these days fine. This points "
        "to a fundamental issue with temporal modeling on early sequences, not a data "
        "preprocessing bug.")

    # =====================================================================
    # SLIDE 5-1: Observation 2 -- Past Signals Are Noise
    # =====================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title(slide, "Obs 2: In Early Days, Past Signals Are Noise")
    add_accent_line(slide, color=ACCENT_ORANGE)

    bullets = [
        "Per-day model wins on 9 of 11 days (Day 6 \u2013 Day 26)",
        ("Advantage ranges from +0.002 to +0.227 balanced accuracy", 1),
        "",
        "Time-series only wins at Day 28 and Day 30",
        ("Advantage: +0.134 to +0.214", 1),
        "",
        "The key insight:",
        ("\u2022 The per-day model ignores history and performs better", 1),
        ("\u2022 If history were informative, the time-series model should dominate", 1),
        ("\u2022 It doesn't \u2192 temporal signal from earlier days is not informative", 1),
        ("\u2022 The past is noise that distorts the time-series model", 1),
        "",
        "This explains Observation 1: noisy past signals \u2192 unstable calibration",
    ]
    add_body_text(slide, bullets, width=Inches(11.5))

    add_speaker_notes(slide,
        "If past images helped, the time-series model should dominate. It loses on 9 of 11 "
        "days. In early days, the organoid's past appearance cannot predict its future quality. "
        "The past is noise. This also explains why the time-series model's threshold calibration "
        "is so unstable -- it's fitting to noise.")

    # =====================================================================
    # SLIDE 5-2: Observation 2 -- Why History Helps at Day 28+
    # =====================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title(slide, "Obs 2 (cont.): Why History Helps Only at Day 28+")
    add_accent_line(slide, color=ACCENT_ORANGE)

    bullets = [
        "By Day 28\u201330, the organoid has shown most of its developmental trajectory",
        ("\u2022 Cumulative visual evidence is finally informative", 1),
        ("\u2022 Quality differentiation happens late \u2014 this makes biological sense", 1),
        "",
        "Earlier days (6\u201326): acceptable and not-acceptable look similar",
        ("\u2022 Morphological ambiguity before differentiation", 1),
        ("\u2022 Adding ambiguous past images only adds noise", 1),
        "",
        "Implication for deployment:",
        ("\u2022 Use per-day model for days 6\u201326", 1),
        ("\u2022 Switch to time-series model for days 28+", 1),
        ("\u2022 Hybrid strategy is optimal", 1),
    ]
    add_body_text(slide, bullets, width=Inches(11.5))

    add_speaker_notes(slide,
        "History becomes signal only near the end of the experiment. This makes biological "
        "sense -- quality differentiation happens late. For deployment, the hybrid strategy "
        "-- per-day early, time-series late -- is optimal.")

    # =====================================================================
    # SLIDE 6-1: Observation 3 -- Small-N Fragility
    # =====================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title(slide, "Obs 3: The Hidden Small-N Fragility Problem")
    add_accent_line(slide, color=ACCENT_GREEN)

    bullets = [
        "5,168 total images sounds large, but...",
        "",
        "The funnel again:",
        ("475 wells \u2192 ~260 labeled \u2192 220 after filtering", 1),
        ("44 in test set, only 9 negatives per day", 1),
        "",
        "Impact of a single misclassification:",
        ("\u2022 1 wrong prediction out of 9 negatives = specificity changes by 11%", 1),
        ("\u2022 2 wrong predictions = specificity changes by 22%", 1),
        "",
        "The 0.227 balanced_acc gap (Day 13, per-day vs effnet_ts)",
        "could be driven by just 2\u20133 samples",
        "",
        "Many performance differences may not be statistically significant.",
    ]
    add_body_text(slide, bullets, width=Inches(11.5))

    add_speaker_notes(slide,
        "5,168 images sounds like a lot, but after filtering and splitting, the test set has "
        "only 9 negative samples per day. A single misclassification flips specificity by 11%. "
        "This means many of the performance differences we've shown may not be statistically "
        "significant.")

    # =====================================================================
    # SLIDE 6-2: Observation 3 -- Implications
    # =====================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title(slide, "Obs 3 (cont.): What This Means for All Our Results")
    add_accent_line(slide, color=ACCENT_GREEN)

    bullets = [
        "Threshold sensitivity (Obs 1) is amplified by small n:",
        ("\u2022 Shifting the threshold moves a few predictions", 1),
        ("\u2022 With 9 negatives, that's a huge percentage swing", 1),
        "",
        "Reproducibility runs (3 runs, fixed seed) produce identical results",
        ("\u2022 The instability is systematic, not random", 1),
        ("\u2022 Same seed \u2192 same result, but the evaluation is still fragile", 1),
        "",
        "What we need:",
        ("\u2022 More labeled negatives (\"Not Acceptable\" organoids)", 1),
        ("\u2022 Cross-validation with confidence intervals", 1),
        ("\u2022 This is the single most impactful improvement", 1),
    ]
    add_body_text(slide, bullets, width=Inches(11.5))

    add_speaker_notes(slide,
        "This is the sobering meta-observation. Before we over-interpret any metric difference, "
        "we need to acknowledge that our evaluation is statistically fragile. More labeled data "
        "-- especially 'not acceptable' -- is the single most impactful improvement we can make.")

    # =====================================================================
    # SLIDE 7-1: Can the Classifier Replace Manual Scoring?
    # =====================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title(slide, "Can the Classifier Replace Manual Scoring?")
    add_accent_line(slide)

    bullets = [
        "Partial yes at late days:",
        ("\u2022 Day 28\u201330: effnet_ts achieves 0.76\u20130.82 balanced accuracy", 1),
        ("\u2022 Useful as a screening tool / pre-filter", 1),
        "",
        "Not yet for early days:",
        ("\u2022 Days 6\u201317: performance 0.50\u20130.67 \u2014 not reliable enough", 1),
        ("\u2022 Cannot replace expert review at early time points", 1),
        "",
        "Cost-saving potential:",
        ("\u2022 Pre-filter at late days: flag clear acceptables automatically", 1),
        ("\u2022 Route ambiguous cases to experts", 1),
        ("\u2022 Could reduce manual review by ~50% at late days", 1),
    ]
    add_body_text(slide, bullets, width=Inches(11.5))

    add_speaker_notes(slide,
        "The classifier can save costs as a pre-filter for late-day assessments, but cannot "
        "fully replace expert review, especially at early time points. At day 28-30 it is "
        "useful as an automated screening tool.")

    # =====================================================================
    # SLIDE 7-2: What Is Wrong with the Image Data
    # =====================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title(slide, "What Is Wrong with the Image Data")
    add_accent_line(slide)

    bullets = [
        "Class imbalance: 72.5% Acceptable",
        ("\u2022 Models default to majority class despite 2.5\u00d7 weighting", 1),
        "",
        "Label propagation assumption:",
        ("\u2022 Labels from Day 28/30 applied backward to earlier days", 1),
        ("\u2022 Assumes quality is determined at seeding \u2014 may not be true", 1),
        "",
        "Small labeled set:",
        ("\u2022 Only 220 usable organoids out of 475", 1),
        ("\u2022 Only 9 negatives in the test set per day", 1),
        "",
        "Morphological ambiguity:",
        ("\u2022 In early days, acceptable and not-acceptable organoids look similar", 1),
        ("\u2022 Not solvable by better models alone", 1),
    ]
    add_body_text(slide, bullets, width=Inches(11.5))

    add_speaker_notes(slide,
        "The data has structural problems that better models alone cannot solve: too few "
        "labeled negatives, a questionable label propagation assumption, and inherent visual "
        "ambiguity in early development.")

    # =====================================================================
    # SLIDE 7-3: Concrete Recommendations
    # =====================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title(slide, "Recommendations")
    add_accent_line(slide)

    bullets = [
        "Short-term:",
        ("\u2022 Deploy per-day model at Day 28\u201330 as automated pre-screening", 1),
        ("\u2022 Hybrid strategy: per-day for early days, time-series for late days", 1),
        "",
        "Medium-term:",
        ("\u2022 Collect more \"Not Acceptable\" labels (target: 100+ negatives in test)", 1),
        ("\u2022 Report all metrics with confidence intervals", 1),
        "",
        "Long-term:",
        ("\u2022 Investigate whether organoid quality is truly fixed at seeding", 1),
        ("\u2022 (the backward label propagation assumption)", 1),
        "",
        "Overlay approach works well \u2192 shape is signal",
        ("\u2022 Simpler shape-based methods could be viable alternatives", 1),
    ]
    add_body_text(slide, bullets, width=Inches(11.5))

    add_speaker_notes(slide,
        "Three concrete next steps: deploy late-day screening now, collect more negative "
        "labels, and question whether backward label propagation is biologically valid. "
        "The fact that overlay images work well suggests simpler shape-based methods could "
        "also be effective.")

    # =====================================================================
    # SLIDE 8-1: Bilinear Combined Model
    # =====================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title(slide, "Bilinear Combined Model: Image + Metabolite")
    add_accent_line(slide)

    bullets = [
        "Tested 7 fusion strategies to add metabolite data into the image classifier",
        ("(concat, gated, FiLM, bilinear, prior, classifier_only, cross_attention)", 1),
        "",
        "Best result: Bilinear fusion",
        ("\u2022 Day 30: 0.87 balanced accuracy (best across all approaches)", 1),
        ("\u2022 Day 28: 0.75 balanced accuracy", 1),
        ("\u2022 Days 8, 13, 20.5: collapses to ~0.50 (same problem as image-only)", 1),
        "",
        "Confirms Observation 2:",
        ("\u2022 Late-day data is where the signal lives \u2014 for both images and metabolites", 1),
        ("\u2022 Metabolite signal is also noisy in early days", 1),
        ("\u2022 No single fusion strategy wins every day", 1),
        "",
        "Next step: Run same 7 strategies with time-series backbone",
    ]
    add_body_text(slide, bullets, width=Inches(11.5))

    add_speaker_notes(slide,
        "We tested 7 ways to combine metabolites with images. The bilinear approach hits "
        "0.87 at day 30 -- our best result. But it doesn't help early days. This confirms "
        "what we learned: late-day data is where the signal lives. The combined approach is "
        "promising but needs further validation with the time-series backbone.")

    # =====================================================================
    # Save
    # =====================================================================
    out_path = "/home/tonyluo/amanda_temporal/image_classifier_presentation.pptx"
    prs.save(out_path)
    print(f"\nPresentation saved to: {out_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
